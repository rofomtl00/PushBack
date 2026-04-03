"""
parser.py — Multi-Format Document Parser
==========================================
Extracts text, tables, and metadata from business documents.
Supports: PDF, DOCX, XLSX, PPTX, CSV, TXT, images.
"""

import os
import csv
import json


def parse_file(filepath: str) -> dict:
    """Parse a file and return structured content."""
    ext = os.path.splitext(filepath)[1].lower()
    name = os.path.basename(filepath)
    size_kb = round(os.path.getsize(filepath) / 1024, 1)

    result = {
        "filename": name,
        "type": ext,
        "size_kb": size_kb,
        "text": "",
        "tables": [],
        "metadata": {},
        "error": None,
    }

    try:
        if ext == ".pdf":
            result.update(_parse_pdf(filepath))
        elif ext == ".docx":
            result.update(_parse_docx(filepath))
        elif ext in (".xlsx", ".xls"):
            result.update(_parse_xlsx(filepath))
        elif ext == ".pptx":
            result.update(_parse_pptx(filepath))
        elif ext == ".csv":
            result.update(_parse_csv(filepath))
        elif ext in (".txt", ".md"):
            with open(filepath, "r", errors="ignore") as f:
                result["text"] = f.read(100000)  # 100KB max
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            result["text"] = f"[Image: {name}]"
            result["metadata"]["type"] = "image"
        else:
            result["text"] = f"[Unsupported file type: {ext}]"
            result["error"] = f"Unsupported format: {ext}"
    except Exception as e:
        result["error"] = str(e)

    return result


def _parse_pdf(filepath):
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return {"text": "", "error": "PyPDF2 not installed"}

    reader = PdfReader(filepath)
    pages = []
    for i, page in enumerate(reader.pages[:50]):  # Max 50 pages
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)

    metadata = {}
    if reader.metadata:
        for key in ("title", "author", "subject", "creator"):
            val = getattr(reader.metadata, key, None)
            if val:
                metadata[key] = str(val)

    return {
        "text": "\n\n--- Page Break ---\n\n".join(pages),
        "metadata": {**metadata, "pages": len(reader.pages)},
    }


def _parse_docx(filepath):
    try:
        from docx import Document
    except ImportError:
        return {"text": "", "error": "python-docx not installed"}

    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)

    # Extract tables
    tables = []
    for table in doc.tables[:10]:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            tables.append(rows)

    return {"text": text, "tables": tables}


def _parse_xlsx(filepath):
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"text": "", "error": "openpyxl not installed"}

    wb = load_workbook(filepath, data_only=True)
    sheets = []
    tables = []

    for ws in wb.worksheets[:10]:
        rows = []
        for row in ws.iter_rows(max_row=200, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            sheets.append({"sheet": ws.title, "rows": len(rows)})
            tables.append(rows)

    # Build text summary
    text_parts = []
    for i, table in enumerate(tables):
        sheet_name = sheets[i]["sheet"] if i < len(sheets) else f"Sheet {i+1}"
        text_parts.append(f"## {sheet_name}")
        for row in table[:100]:  # Max 100 rows per sheet
            text_parts.append(" | ".join(row))

    return {
        "text": "\n".join(text_parts),
        "tables": tables,
        "metadata": {"sheets": [s["sheet"] for s in sheets]},
    }


def _parse_pptx(filepath):
    try:
        from pptx import Presentation
    except ImportError:
        return {"text": "", "error": "python-pptx not installed"}

    prs = Presentation(filepath)
    slides = []

    for i, slide in enumerate(prs.slides[:50]):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            slides.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))

    return {
        "text": "\n\n".join(slides),
        "metadata": {"slides": len(prs.slides)},
    }


def _parse_csv(filepath):
    with open(filepath, "r", errors="ignore") as f:
        reader = csv.reader(f)
        rows = []
        for i, row in enumerate(reader):
            if i >= 200:
                break
            rows.append(row)

    text_parts = []
    for row in rows:
        text_parts.append(" | ".join(row))

    return {"text": "\n".join(text_parts), "tables": [rows]}


def parse_folder(filepaths: list) -> dict:
    """Parse multiple files and return combined context."""
    results = []
    total_text = []

    for fp in filepaths:
        result = parse_file(fp)
        results.append(result)
        if result["text"]:
            total_text.append(f"=== {result['filename']} ===\n{result['text']}")

    combined_text = "\n\n".join(total_text)
    # Truncate to ~100K chars to fit in AI context
    if len(combined_text) > 100000:
        combined_text = combined_text[:100000] + "\n\n[Content truncated — too large]"

    return {
        "files": results,
        "file_count": len(results),
        "total_chars": len(combined_text),
        "combined_text": combined_text,
        "errors": [r for r in results if r.get("error")],
    }
